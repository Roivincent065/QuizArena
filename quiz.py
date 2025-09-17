import streamlit as st
import base64
import json
import time
import random
from datetime import datetime
import groq
import PyPDF2
import docx
from pptx import Presentation
import io
import pandas as pd
import os
import hashlib
import threading

# Set up the page
st.set_page_config(
    page_title="QuizArena - Gamified Learning",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize Groq client
if "groq_client" not in st.session_state:
    if "GROQ_API_KEY" in st.secrets and st.secrets["GROQ_API_KEY"]:
        try:
            st.session_state.groq_client = groq.Client(api_key=st.secrets["GROQ_API_KEY"])
        except Exception as e:
            st.error(f"Error initializing Groq client: {e}")
            st.session_state.groq_client = None
    else:
        st.session_state.groq_client = None

# Initialize session state variables
if "is_logged_in" not in st.session_state:
    st.session_state.is_logged_in = False
if "user_id" not in st.session_state:
    st.session_state.user_id = None
if "username" not in st.session_state:
    st.session_state.username = ""
if "avatar" not in st.session_state:
    st.session_state.avatar = "🧠"
if "current_page" not in st.session_state:
    st.session_state.current_page = "login"
if "lobbies" not in st.session_state:
    st.session_state.lobbies = {}
if "current_lobby" not in st.session_state:
    st.session_state.current_lobby = None
if "quiz_data" not in st.session_state:
    st.session_state.quiz_data = None
if "game_started" not in st.session_state:
    st.session_state.game_started = False
if "user_answers" not in st.session_state:
    st.session_state.user_answers = {}
if "leaderboard" not in st.session_state:
    st.session_state.leaderboard = {}
if "user_score" not in st.session_state:
    st.session_state.user_score = 0
if "streak" not in st.session_state:
    st.session_state.streak = 0
if "trivia_data" not in st.session_state:
    st.session_state.trivia_data = None
if "trivia_categories" not in st.session_state:
    st.session_state.trivia_categories = []
if "question_start_time" not in st.session_state:
    st.session_state.question_start_time = None
if "timer_active" not in st.session_state:
    st.session_state.timer_active = False
if "selected_answer" not in st.session_state:
    st.session_state.selected_answer = None
if "answer_submitted" not in st.session_state:
    st.session_state.answer_submitted = False
if "prev_page" not in st.session_state:
    st.session_state.prev_page = "home"
if "chat_messages" not in st.session_state:
    st.session_state.chat_messages = {}
if "current_question" not in st.session_state:
    st.session_state.current_question = 0
if "start_time" not in st.session_state:
    st.session_state.start_time = time.time()

# Kahoot-like colors for options
OPTION_COLORS = ["#FF2B2B", "#1E88E5", "#FFC107", "#4CAF50"]
OPTION_LABELS = ["🟥", "🟦", "🟨", "🟩"]
EMOJI_AVATARS = ["🧠", "🚀", "💡", "📚", "🎓", "🌟", "🤓", "😎", "🧐", "🤔"]

# --- Database Functions (JSON file) ---
USERS_DB = "users.json"
LOBBIES_DB = "lobbies.json"

def load_users():
    if not os.path.exists(USERS_DB):
        return {}
    with open(USERS_DB, "r") as f:
        return json.load(f)

def save_users(users):
    with open(USERS_DB, "w") as f:
        json.dump(users, f, indent=4)

def load_lobbies():
    if not os.path.exists(LOBBIES_DB):
        return {}
    with open(LOBBIES_DB, "r") as f:
        return json.load(f)

def save_lobbies(lobbies):
    with open(LOBBIES_DB, "w") as f:
        json.dump(lobbies, f, indent=4)

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# --- Utility Functions ---
def set_page(page, prev_page=None):
    if prev_page:
        st.session_state.prev_page = prev_page
    st.session_state.current_page = page
    st.rerun()

# Function to load general knowledge data
def load_trivia_data():
    try:
        data_path = "data/general_knowledge_qa.csv"
        if os.path.exists(data_path):
            df = pd.read_csv(data_path)
            st.session_state.trivia_data = df
            if 'category' in df.columns:
                st.session_state.trivia_categories = df['category'].unique().tolist()
            return True
        else:
            st.session_state.trivia_data = pd.DataFrame({
                'question': ['What is the capital of France?','Which planet is known as the Red Planet?','Who painted the Mona Lisa?','What is the largest mammal in the world?','In which year did World War II end?','What is the chemical symbol for gold?','Who wrote "Romeo and Juliet"?','What is the largest ocean on Earth?','How many elements are in the periodic table?','What is the tallest mountain in the world?'],
                'answer': ['Paris','Mars','Leonardo da Vinci','Blue Whale','1945','Au','William Shakespeare','Pacific Ocean','118','Mount Everest'],
                'options': ['Paris|London|Berlin|Madrid','Mars|Venus|Jupiter|Saturn','Leonardo da Vinci|Pablo Picasso|Vincent van Gogh|Michelangelo','Blue Whale|Elephant|Giraffe|Hippopotamus','1945|1918|1939|1941','Au|Ag|Fe|Cu','William Shakespeare|Charles Dickens|Jane Austen|Mark Twain','Pacific Ocean|Atlantic Ocean|Indian Ocean|Arctic Ocean','118|92|108|132','Mount Everest|K2|Kilimanjaro|Matterhorn'],
                'category': ['Geography','Science','Art','Science','History','Science','Literature','Geography','Science','Geography'],
                'difficulty': ['Easy','Easy','Medium','Medium','Hard','Medium','Easy','Easy','Hard','Medium']
            })
            st.session_state.trivia_categories = ['Geography', 'Science', 'Art', 'History', 'Literature']
            return True
    except Exception as e:
        st.error(f"Error loading trivia data: {str(e)}")
        return False

# Function to generate trivia quiz from the dataset
def generate_trivia_quiz(category=None, difficulty=None, num_questions=5):
    if st.session_state.trivia_data is None:
        if not load_trivia_data():
            return None
    
    df = st.session_state.trivia_data
    if category and category != "All" and 'category' in df.columns:
        df = df[df['category'] == category]
    if difficulty and difficulty != "All" and 'difficulty' in df.columns:
        df = df[df['difficulty'] == difficulty]
    
    if len(df) > num_questions:
        df = df.sample(n=num_questions)
    
    quiz_data = {
        "quiz_title": f"General Knowledge Trivia - {category if category else 'All Categories'}",
        "questions": []
    }
    
    for _, row in df.iterrows():
        question = {
            "question": row['question'],
            "correct_answer": row['answer'],
            "question_type": "mcq"
        }
        if 'options' in row and pd.notna(row['options']) and len(row['options'].split('|')) == 4:
            options = row['options'].split('|')
            random.shuffle(options)
            question["options"] = options
        else:
            all_answers = st.session_state.trivia_data['answer'].tolist()
            correct_answer = row['answer']
            incorrect_answers = [ans for ans in all_answers if ans != correct_answer]
            
            if len(incorrect_answers) >= 3:
                options = random.sample(incorrect_answers, 3)
                options.append(correct_answer)
                random.shuffle(options)
            else:
                options = [correct_answer] + ["N/A"] * 3
                random.shuffle(options)
            question["options"] = options
        quiz_data["questions"].append(question)
    return quiz_data

# Function to extract text from different file types
def extract_text_from_file(file):
    file_type = file.type
    text = ""
    try:
        if file_type == "text/plain":
            text = str(file.read(), "utf-8")
        elif file_type == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(io.BytesIO(file.read()))
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return ""

# Function to generate quiz using Groq API
def generate_quiz(text, game_mode, num_questions=5):
    if not st.session_state.groq_client:
        st.error("Groq API key not configured or is invalid. Please check your secrets.toml file.")
        st.info("You can get a Groq API key from https://console.groq.com/keys")
        return None
    
    if game_mode == "Multiple Choice":
        prompt = f"""Create a multiple choice quiz based on the following text. Generate {num_questions} questions with 4 options each and indicate the correct answer.
        Format your response as JSON with the following structure:
        {{"quiz_title": "Generated Quiz","questions": [{{"question": "Question text","options": ["Option1", "Option2", "Option3", "Option4"],"correct_answer": "Option1","question_type": "mcq"}}]}}
        Text: {text}
        """
    elif game_mode == "True or False":
        prompt = f"""Create a True or False quiz based on the following text. Generate {num_questions} questions and indicate whether each statement is true or false.
        Format your response as JSON with the following structure:
        {{"quiz_title": "Generated Quiz","questions": [{{"question": "Statement text","correct_answer": "True","question_type": "true_false"}}]}}
        Text: {text}
        """
    elif game_mode == "Identification":
        prompt = f"""Create an identification quiz (fill-in-the-blank) based on the following text. Generate {num_questions} questions with clear answers.
        Format your response as JSON with the following structure:
        {{"quiz_title": "Generated Quiz","questions": [{{"question": "Question text with _____ for blank","correct_answer": "Answer","question_type": "identification"}}]}}
        Text: {text}
        """
    elif game_mode == "Enumeration":
        prompt = f"""Create an enumeration quiz based on the following text. Generate {num_questions} questions that ask for lists of items, with each item separated by commas in the correct answer.
        Format your response as JSON with the following structure:
        {{"quiz_title": "Generated Quiz","questions": [{{"question": "Question text asking for a list","correct_answer": "Item1, Item2, Item3","question_type": "enumeration"}}]}}
        Text: {text}
        """
    elif game_mode == "Mix Mode":
        prompt = f"""Create a mixed format quiz based on the following text. Generate {num_questions} questions with a variety of types (multiple choice, true/false, identification, enumeration).
        Format your response as JSON with the following structure:
        {{"quiz_title": "Generated Quiz","questions": [{{"question": "Question text","options": ["Option1", "Option2", "Option3", "Option4"] (only for multiple choice),"correct_answer": "Answer","question_type": "mcq/true_false/identification/enumeration"}}]}}
        Text: {text}
        """
    
    try:
        chat_completion = st.session_state.groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.1-8b-instant",
            temperature=0.7,
            max_tokens=4000
        )
        response = chat_completion.choices[0].message.content
        json_start = response.find('{')
        json_end = response.rfind('}') + 1
        json_str = response[json_start:json_end]
        quiz_data = json.loads(json_str)
        return quiz_data
    except Exception as e:
        st.error(f"Error generating quiz: {str(e)}")
        return None

# Function to create a new lobby
def create_lobby(lobby_name, lobby_type, max_players=10):
    lobbies = load_lobbies()
    lobby_id = f"L{random.randint(10000, 99999)}"
    lobbies[lobby_id] = {
        "id": lobby_id,
        "name": lobby_name,
        "type": lobby_type,
        "max_players": max_players,
        "players": [st.session_state.user_id],
        "player_names": [st.session_state.username],
        "host": st.session_state.user_id,
        "status": "waiting",
        "quiz_data": None,
        "scores": {st.session_state.user_id: 0},
        "start_time": None,
        "chat_messages": [],
        "votes_to_start": {}
    }
    save_lobbies(lobbies)
    return lobby_id

# Function to join a lobby
def join_lobby(lobby_id):
    lobbies = load_lobbies()
    if lobby_id in lobbies:
        lobby = lobbies[lobby_id]
        if len(lobby["players"]) < lobby["max_players"] and st.session_state.user_id not in lobby["players"]:
            lobby["players"].append(st.session_state.user_id)
            lobby["player_names"].append(st.session_state.username)
            lobby["scores"][st.session_state.user_id] = 0
            save_lobbies(lobbies)
            return True
    return False

# Function to start the game in a lobby (after voting)
def start_game(lobby_id):
    lobbies = load_lobbies()
    if lobby_id in lobbies:
        lobbies[lobby_id]["status"] = "playing"
        st.session_state.quiz_data = lobbies[lobby_id]["quiz_data"]
        lobbies[lobby_id]["current_question"] = 0
        lobbies[lobby_id]["start_time"] = time.time()
        lobbies[lobby_id]["question_start_time"] = time.time()
        save_lobbies(lobbies)
        st.session_state.game_started = True
        set_page("playing", "lobby_page") # Correctly set prev_page
        return True
    return False

# Function to check answer
def check_answer(question, user_answer, qtype):
    if not user_answer:
        return False
    correct_answer = question.get("correct_answer")
    if not correct_answer:
        return False
    if qtype in ["mcq", "true_false", "identification"]:
        return str(user_answer).strip().lower() == str(correct_answer).strip().lower()
    elif qtype == "enumeration":
        correct_answers = [a.strip().lower() for a in str(correct_answer).split(",")]
        return str(user_answer).strip().lower() in correct_answers
    elif qtype == "essay":
        return str(user_answer).strip() != ""
    return False

# Function to calculate score based on time and accuracy
def calculate_score(time_taken, is_correct, question_type, accuracy=1.0):
    base_score = 100
    time_bonus = max(0, 5 - time_taken) * 20
    if is_correct:
        if question_type == "enumeration":
            return int(base_score * accuracy + time_bonus)
        return int(base_score + time_bonus)
    return 0

# --- Page Functions ---

# Login/Registration Page
def login_page():
    st.markdown("""
    <style>
    .login-container {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-top: 5rem;
        font:Old English Text MT;
    }
    .stTextInput>div>div>input {
        color: black;
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown('<div class="login-container"><h1>🎓 QuizArena - Login</h1></div>', unsafe_allow_html=True)
    
    tab1, tab2 = st.tabs(["🔒 Login", "✍️ Register"])

    with tab1:
        st.subheader("Existing User Login")
        login_username = st.text_input("Nickname", key="login_username")
        login_password = st.text_input("Password", type="password", key="login_password")
        if st.button("Login", use_container_width=True):
            users = load_users()
            if login_username in users and users[login_username]["password"] == hash_password(login_password):
                st.session_state.is_logged_in = True
                st.session_state.username = login_username
                st.session_state.user_id = users[login_username]["user_id"]
                st.session_state.avatar = users[login_username].get("avatar", "🧠")
                st.success(f"Welcome back, {login_username}!")
                set_page("home")
            else:
                st.error("Invalid nickname or password.")

    with tab2:
        st.subheader("New User Registration")
        reg_username = st.text_input("Choose a Nickname", key="reg_username")
        reg_password = st.text_input("Create a Password", type="password", key="reg_password")
        reg_avatar = st.selectbox("Choose your Avatar", EMOJI_AVATARS)
        if st.button("Register", use_container_width=True):
            users = load_users()
            if reg_username in users:
                st.error("This nickname is already taken.")
            elif len(reg_password) < 4:
                st.error("Password must be at least 4 characters long.")
            else:
                new_user_id = f"user_{random.randint(10000, 99999)}"
                users[reg_username] = {
                    "user_id": new_user_id,
                    "password": hash_password(reg_password),
                    "avatar": reg_avatar,
                    "score": 0,
                    "quizzes_completed": 0
                }
                save_users(users)
                st.success("Registration successful! Please log in.")
                
# Edit Profile Page
def edit_profile_page():
    st.title("👤 Edit Profile")
    st.write("Update your nickname or avatar.")

    new_username = st.text_input("New Nickname", value=st.session_state.username)
    new_avatar = st.selectbox("Choose a new Avatar", EMOJI_AVATARS, index=EMOJI_AVATARS.index(st.session_state.avatar))

    if st.button("Save Changes", type="primary"):
        users = load_users()
        # Handle nickname change
        if new_username != st.session_state.username:
            if new_username in users:
                st.error("This nickname is already in use.")
            else:
                # Update users database with new nickname
                user_data = users.pop(st.session_state.username)
                user_data["username"] = new_username
                user_data["avatar"] = new_avatar
                users[new_username] = user_data
                st.session_state.username = new_username
                st.session_state.avatar = new_avatar
                save_users(users)
                st.success("Profile updated successfully!")
        else:
            # Only update avatar
            users[st.session_state.username]["avatar"] = new_avatar
            st.session_state.avatar = new_avatar
            save_users(users)
            st.success("Avatar updated successfully!")
    
    if st.button("← Go Back"):
        set_page("home")

# Home page
def home_page():
    st.markdown("""
    <style>
    .main-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<div class="main-header"><h1>🎓 QuizArena - Gamified Learning</h1></div>', unsafe_allow_html=True)
    
    st.write(f"Welcome, **{st.session_state.username}**! 👋")
    
    col1, col2, col3, col4, col5 = st.columns(5)
    
    with col1:
        if st.button("📚 Quiz Lobby", use_container_width=True, type="primary"):
            set_page("exam_prep", "home")
    with col2:
        if st.button("🎯 General Knowledge Trivia", use_container_width=True, type="primary"):
            set_page("trivia", "home")
    with col3:
        if st.button("🏆 Leaderboards", use_container_width=True, type="primary"):
            set_page("leaderboards", "home")
    with col4:
        if st.button("🧠 Mindfulness Breaks", use_container_width=True, type="primary"):
            set_page("mindfulness", "home")
    with col5:
        if st.button("⚙️ Edit Profile", use_container_width=True, type="secondary"):
            set_page("edit_profile", "home")

    st.markdown("---")
    st.subheader("🎮 How it works:")
    st.markdown("""
    <div style="background-color: #f0f2f6; padding: 20px; border-radius: 10px; border-left: 5px solid #667eea;">
    1. **Create or join** a study lobby with your classmates<br>
    2. **Upload study materials** (PDF, PPTX, DOCX, TXT)<br>
    3. **AI generates quizzes** based on your materials<br>
    4. **Compete in real-time** with various game modes<br>
    5. **Earn points** for correct answers and speed<br>
    6. **Track your progress** on leaderboards
    </div>
    """, unsafe_allow_html=True)

# Quiz Lobby Page
def exam_prep_page():
    st.title("📚 Quiz Lobby")
    if st.button("← Go Back"):
        set_page("home")
    
    st.markdown("""
    <style>
    .lobby-card {
        background: linear-gradient(135deg, #ff6b6b 0%, #ee5a24 100%);
        padding: 1.5rem;
        border-radius: 15px;
        color: white;
        margin-bottom: 1rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Load lobbies at the start of the page function
    st.session_state.lobbies = load_lobbies()

    tab1, tab2 = st.tabs(["🎪 Create Lobby", "🚪 Join Lobby"])
    
    with tab1:
        st.subheader("Create a New Study Lobby")
        lobby_name = st.text_input("Lobby Name", value=f"{st.session_state.username}'s Study Group")
        lobby_type = st.selectbox("Lobby Type", ["Private", "Public"])
        max_players = st.slider("Maximum Players", 2, 20, 10)
        
        if st.button("🎉 Create Lobby", type="primary"):
            lobby_id = create_lobby(lobby_name, lobby_type, max_players)
            st.session_state.current_lobby = lobby_id
            st.success(f"Lobby created! Your lobby code is: **{lobby_id}**")
            set_page("lobby_page", "exam_prep") # Correctly set prev_page
            st.rerun()
    
    with tab2:
        st.subheader("Join an Existing Lobby")
        lobby_id = st.text_input("Enter Lobby Code")
        
        if st.button("🎯 Join Lobby", type="primary"):
            if join_lobby(lobby_id):
                st.session_state.current_lobby = lobby_id
                st.success("Joined lobby successfully! 🎉")
                set_page("lobby_page", "exam_prep") # Correctly set prev_page
                st.rerun()
            else:
                st.error("Could not join lobby. It may be full or doesn't exist.")

# New Lobby Page
def lobby_page():
    st.title("🎪 Lobby")
    
    st.session_state.lobbies = load_lobbies()
    lobby = st.session_state.lobbies.get(st.session_state.current_lobby)

    if not lobby:
        st.error("Lobby not found.")
        st.session_state.current_lobby = None
        set_page("exam_prep")
        st.rerun()
    
    # Check if the game has started
    if lobby["status"] == "playing":
        st.info("The host has started the game!")
        set_page("playing", "lobby_page")
        st.rerun()

    if st.button("← Leave Lobby"):
        st.session_state.current_lobby = None
        set_page("exam_prep")
        st.rerun()

    st.markdown(f'### Lobby: {lobby["name"]} ({lobby["id"]})')
    st.write(f"**Players:** {', '.join(lobby['player_names'])}")

    col1, col2 = st.columns([2, 1])

    with col1:
        st.subheader("Lobby Actions")
        if lobby["host"] == st.session_state.user_id:
            st.markdown("---")
            st.subheader("📁 Upload Materials & Generate Quiz")
            uploaded_file = st.file_uploader("Choose a file", type=["pdf", "pptx", "docx", "txt"])
            game_mode = st.selectbox("Select Game Mode", 
                                    ["Multiple Choice", "True or False", "Identification", "Enumeration", "Mix Mode"])
            num_questions = st.slider("Number of Questions", 5, 20, 10)
            
            if uploaded_file and st.button("⚡ Generate Quiz", type="primary"):
                with st.spinner("Extracting text and generating quiz..."):
                    text = extract_text_from_file(uploaded_file)
                    if text:
                        quiz_data = generate_quiz(text, game_mode, num_questions)
                        if quiz_data:
                            lobby["quiz_data"] = quiz_data
                            save_lobbies(st.session_state.lobbies)
                            st.success("Quiz generated successfully! 🎯")
                            st.rerun()
                        else:
                            st.error("Failed to generate quiz.")
                    else:
                        st.error("Could not extract text from the file.")
            
            if lobby["quiz_data"] and lobby["status"] == "waiting":
                if st.button("🚀 Start Game", type="primary"):
                    st.success("Starting the game...")
                    start_game(st.session_state.current_lobby)
        else:
            st.info("Waiting for the host to upload materials and start the game.")

    with col2:
        st.subheader("💬 Lobby Chat")
        
        # Display chat messages
        chat_placeholder = st.container()
        with chat_placeholder:
            for message in lobby.get("chat_messages", []):
                st.write(f"**{message['username']}:** {message['message']}")
        
        # Input for new message
        chat_input = st.text_input("Type your message here...", key="chat_input")
        if st.button("Send", use_container_width=True) and chat_input:
            new_message = {
                "username": st.session_state.username,
                "message": chat_input,
                "timestamp": datetime.now().isoformat()
            }
            lobby["chat_messages"].append(new_message)
            save_lobbies(st.session_state.lobbies)
            st.rerun()

# Functions for color logic
def get_random_color():
    return f"#{random.randint(0, 0xFFFFFF):06x}"

def get_text_color(hex_color):
    hex_color = hex_color.lstrip('#')
    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    luminance = (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) / 255
    return "white" if luminance < 0.5 else "black"

# Function to play the game
def play_game(quiz_data, current_idx, question_start_time):
    st.markdown("""
    <style>
    .question-container {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 20px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    .timer-container {
        background: #ff4757;
        padding: 1rem;
        border-radius: 50%;
        width: 80px;
        height: 80px;
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: bold;
        font-size: 24px;
        margin: 0 auto;
        box-shadow: 0 4px 8px rgba(0,0,0,0.2);
    }
    .stButton>button {
        background-color: var(--button-color) !important;
        padding: 20px !important;
        border-radius: 15px !important;
        margin: 10px 0 !important;
        color: var(--text-color) !important;
        font-weight: bold !important;
        font-size: 18px !important;
        text-align: center !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.2) !important;
        transition: transform 0.2s !important;
        border: none !important;
        cursor: pointer !important;
        width: 100% !important;
        height: 100px !important; /* Fixed height for uniform size */
        display: flex !important;
        flex-direction: column !important;
        justify-content: center !important;
    }
    </style>
    """, unsafe_allow_html=True)

    questions = quiz_data["questions"]
    
    if current_idx < len(questions):
        question = questions[current_idx]
        
        # Set dynamic time limit based on question type
        if question["question_type"] in ["identification", "enumeration"]:
            answer_length = len(question.get("correct_answer", ""))
            time_limit = min(15, max(10, 10 + answer_length // 5)) # 10-15 seconds
        else:
            time_limit = 10
        
        # Display the question and timer
        st.markdown(f'<div class="question-container"><h2>Question {current_idx + 1} of {len(questions)}</h2><h3>{question["question"]}</h3></div>', unsafe_allow_html=True)
        timer_placeholder = st.empty()

        elapsed_time = time.time() - question_start_time
        time_remaining = int(max(0, time_limit - elapsed_time))
        
        with timer_placeholder.container():
            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                st.markdown(f'<div class="timer-container">{time_remaining}s</div>', unsafe_allow_html=True)

        # Check if timer has run out
        if time_remaining <= 0:
            if not st.session_state.answer_submitted:
                st.session_state.selected_answer = "Time's up!"
            
            time_taken = int(time.time() - question_start_time)
            is_correct = check_answer(question, st.session_state.selected_answer, question["question_type"])
            accuracy = 1.0
            score = calculate_score(time_taken, is_correct, question["question_type"], accuracy)
            
            st.session_state.user_score += score
            if st.session_state.current_lobby:
                lobbies = load_lobbies()
                lobbies[st.session_state.current_lobby]["scores"][st.session_state.user_id] += score
                save_lobbies(lobbies)
            if is_correct:
                st.session_state.streak += 1
            else:
                st.session_state.streak = 0
            
            st.session_state.user_answers[current_idx] = {
                "user_answer": st.session_state.selected_answer,
                "correct_answer": question["correct_answer"],
                "is_correct": is_correct,
                "score": score,
                "time_taken": time_taken
            }
            
            st.markdown("---")
            if is_correct:
                st.success("✅ Correct!")
            else:
                st.error("❌ Incorrect!")
            
            st.info(f"The correct answer was: **{question['correct_answer']}**")
            st.info(f"You earned: **{int(score)}** points!")
            
            time.sleep(3)
            
            if st.session_state.current_lobby:
                lobbies = load_lobbies()
                lobbies[st.session_state.current_lobby]["current_question"] += 1
                lobbies[st.session_state.current_lobby]["question_start_time"] = time.time()
                save_lobbies(lobbies)
            else:
                st.session_state.current_question += 1
                st.session_state.question_start_time = time.time()

            st.session_state.selected_answer = None
            st.session_state.answer_submitted = False
            st.rerun()

        # Display options based on question type
        if not st.session_state.answer_submitted:
            if question["question_type"] == "mcq":
                options = question.get("options", [])
                cols = st.columns(2)
                for i, option in enumerate(options):
                    with cols[i % 2]:
                        random_color = OPTION_COLORS[i % 4]
                        text_color = get_text_color(random_color)
                        button_style = f"background-color: {random_color}; color: {text_color};"
                        
                        if st.button(option, key=f"q{current_idx}_{i}", use_container_width=True):
                            st.session_state.selected_answer = option
                            st.session_state.answer_submitted = True
                            st.rerun()

            elif question["question_type"] == "true_false":
                options = ["True", "False"]
                col1, col2 = st.columns(2)
                for i, option in enumerate(options):
                    with locals()[f"col{i+1}"]:
                        random_color = OPTION_COLORS[i % 2]
                        text_color = get_text_color(random_color)
                        button_style = f"background-color: {random_color}; color: {text_color};"
                        
                        if st.button(option, key=f"q{current_idx}_{i}", use_container_width=True):
                            st.session_state.selected_answer = option
                            st.session_state.answer_submitted = True
                            st.rerun()
            else:
                user_answer = st.text_input("Your answer:", key=f"q{current_idx}")
                if st.button("Submit Answer", key=f"submit_{current_idx}"):
                    st.session_state.selected_answer = user_answer
                    st.session_state.answer_submitted = True
                    st.rerun()
        
        # Rerun to update timer if time is still left
        if time_remaining > 0:
            time.sleep(1)
            st.rerun()
            
    else:
        # Quiz completed
        st.balloons()
        st.markdown('<div class="question-container"><h2>🎉 Quiz Completed!</h2></div>', unsafe_allow_html=True)
        
        # Update user's global score
        users = load_users()
        if st.session_state.username in users:
            users[st.session_state.username]["score"] += st.session_state.user_score
            users[st.session_state.username]["quizzes_completed"] += 1
            save_users(users)

        # Display Match Leaderboard
        st.subheader("🏆 Match Leaderboard")
        match_scores = []
        if st.session_state.current_lobby:
            lobbies = load_lobbies()
            lobby = lobbies[st.session_state.current_lobby]
            for user_id, score in lobby["scores"].items():
                username = next((name for p_id, name in zip(lobby["players"], lobby["player_names"]) if p_id == user_id), "Unknown")
                match_scores.append({"Username": username, "Score": int(score)})
        else:
            match_scores.append({"Username": st.session_state.username, "Score": int(st.session_state.user_score)})

        match_df = pd.DataFrame(match_scores)
        match_df = match_df.sort_values(by="Score", ascending=False).reset_index(drop=True)
        match_df.index = match_df.index + 1
        match_df.insert(0, 'Rank', match_df.index)
        st.dataframe(match_df, use_container_width=True, hide_index=True)
        
        # Action buttons
        if st.button("🔄 Play Again", type="primary"):
            st.session_state.current_question = 0
            st.session_state.user_answers = {}
            st.session_state.user_score = 0
            st.session_state.game_started = False
            st.session_state.selected_answer = None
            st.session_state.answer_submitted = False
            st.session_state.streak = 0
            st.session_state.question_start_time = time.time()  # Reset timer for next game
            st.session_state.current_page = st.session_state.prev_page
            st.rerun()
            
        if st.session_state.prev_page == "lobby_page":
            if st.button("← Go Back to Quiz Lobby"):
                st.session_state.current_question = 0
                st.session_state.user_answers = {}
                st.session_state.user_score = 0
                st.session_state.game_started = False
                st.session_state.selected_answer = None
                st.session_state.answer_submitted = False
                st.session_state.streak = 0
                st.session_state.question_start_time = time.time()  # Reset timer for next game
                set_page("lobby_page")
        else:
            if st.button("← Go Back to Trivia Page"):
                st.session_state.current_question = 0
                st.session_state.user_answers = {}
                st.session_state.user_score = 0
                st.session_state.game_started = False
                st.session_state.selected_answer = None
                st.session_state.answer_submitted = False
                st.session_state.streak = 0
                st.session_state.question_start_time = time.time()  # Reset timer for next game
                set_page("trivia")

# Trivia page
def trivia_page():
    st.markdown("""
    <style>
    .trivia-header {
        background: linear-gradient(135deg, #ffd89b 0%, #19547b 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<div class="trivia-header"><h1>🎯 General Knowledge Trivia</h1></div>', unsafe_allow_html=True)
    if st.button("← Go Back"):
        set_page("home")
    
    if st.session_state.trivia_data is None:
        load_trivia_data()
    
    st.subheader("Test your general knowledge! 🧠")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.session_state.trivia_categories:
            category = st.selectbox("Select Category", ["All"] + st.session_state.trivia_categories)
        else:
            category = "All"
        
        difficulty = st.selectbox("Select Difficulty", ["All", "Easy", "Medium", "Hard"])
    
    with col2:
        num_questions = st.slider("Number of Questions", 5, 20, 10)
        
        if st.button("🚀 Start Trivia Quiz", type="primary"):
            quiz_data = generate_trivia_quiz(category, difficulty, num_questions)
            if quiz_data:
                st.session_state.quiz_data = quiz_data
                st.session_state.current_question = 0
                st.session_state.user_answers = {}
                st.session_state.user_score = 0
                st.session_state.game_started = True
                st.session_state.current_lobby = None # Ensure no lobby is tied to this game
                st.session_state.question_start_time = time.time()  # Reset timer for new game
                set_page("playing", "trivia")
            else:
                st.error("Could not generate trivia quiz. Please try again.")

# Playing page (for both exam prep and trivia)
def playing_page():
    if st.session_state.current_lobby:
        lobbies = load_lobbies()
        lobby = lobbies[st.session_state.current_lobby]
        quiz_data = lobby["quiz_data"]
        current_idx = lobby["current_question"]
        question_start_time = lobby["question_start_time"]
    else:
        quiz_data = st.session_state.quiz_data
        current_idx = st.session_state.current_question
        question_start_time = st.session_state.question_start_time

    if quiz_data:
        play_game(quiz_data, current_idx, question_start_time)
    else:
        st.error("No quiz data found. Please go back and generate a quiz first.")
        if st.button("← Go Back"):
            st.session_state.game_started = False
            set_page(st.session_state.prev_page)

# Leaderboards page
def leaderboards_page():
    st.markdown("""
    <style>
    .leaderboard-header {
        background: linear-gradient(135deg, #f46b45 0%, #eea849 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<div class="leaderboard-header"><h1>🏆 Leaderboards</h1></div>', unsafe_allow_html=True)
    if st.button("← Go Back"):
        set_page("home")

    st.subheader("🌍 Global Leaderboard")
    users = load_users()
    if users:
        leaderboard_data = []
        for username, data in users.items():
            leaderboard_data.append({
                "Rank": "-",
                "Avatar": data.get("avatar", "🧠"),
                "Username": username,
                "Total Score": data["score"],
                "Quizzes Completed": data["quizzes_completed"]
            })
        
        df = pd.DataFrame(leaderboard_data)
        df = df.sort_values(by="Total Score", ascending=False).reset_index(drop=True)
        df.index = df.index + 1
        df["Rank"] = df.index
        df = df[["Rank", "Avatar", "Username", "Total Score", "Quizzes Completed"]]
        
        st.dataframe(df, hide_index=True)

    else:
        st.info("No leaderboard data yet. Complete some quizzes to appear here! 🎯")

# Mindfulness page
def mindfulness_page():
    st.markdown("""
    <style>
    .mindfulness-header {
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<div class="mindfulness-header"><h1>🧠 Mindfulness Breaks</h1></div>', unsafe_allow_html=True)
    if st.button("← Go Back"):
        set_page("home")
    
    st.write("Take a break and relax with these mindfulness exercises. 🌿")
    
    tab1, tab2 = st.tabs(["🌬️ Breathing Exercise", "🎯 Focus Game"])
    
    with tab1:
        st.subheader("Deep Breathing Exercise")
        st.write("Follow the animation and breathe in and out slowly.")
        
        breath_duration = st.slider("Breath Duration (seconds)", 3, 10, 5)
        
        if st.button("Start Breathing Exercise", type="primary"):
            breathing_placeholder = st.empty()
            for i in range(3):
                breathing_placeholder.info("🌬️ Breathe IN...")
                time.sleep(breath_duration)
                breathing_placeholder.info("⏸️ Hold...")
                time.sleep(2)
                breathing_placeholder.info("💨 Breathe OUT...")
                time.sleep(breath_duration)
            
            breathing_placeholder.success("✅ Exercise completed! Feel more relaxed? 😊")
    
    with tab2:
        st.subheader("Focus Game")
        st.write("Watch the circle and try to keep it centered.")
        
        if st.button("Start Focus Game", type="primary"):
            focus_placeholder = st.empty()
            for i in range(10):
                focus_placeholder.markdown(
                    f"<div style='text-align: center; font-size: 50px;'>◉</div>", 
                    unsafe_allow_html=True
                )
                time.sleep(1)
            
            focus_placeholder.success("✅ Focus exercise completed! 🎯")

# Main app routing
def main():
    # Sidebar
    with st.sidebar:
        st.markdown("""
        <style>
        .sidebar-header {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 1rem;
            border-radius: 10px;
            color: white;
            text-align: center;
            margin-bottom: 1rem;
        }
        </style>
        """, unsafe_allow_html=True)
        
        st.markdown('<div class="sidebar-header"><h2>🎓 QuizArena</h2></div>', unsafe_allow_html=True)
        
        if st.session_state.is_logged_in:
            st.write(f"{st.session_state.avatar} User: **{st.session_state.username}**")
            st.write(f"⭐ Score: **{int(st.session_state.user_score)}**")
            st.write(f"🔥 Streak: **{st.session_state.streak}**")
            
            st.markdown("---")
            if st.button("🏠 Home", use_container_width=True):
                set_page("home")
            if st.button("🚪 Logout", use_container_width=True):
                st.session_state.clear()
                st.rerun()

        st.markdown("---")
        st.write("ℹ️ About QuizArena")
        st.caption("A gamified learning platform that makes studying fun and collaborative! 🎯")
    
    # Page routing
    if st.session_state.is_logged_in:
        if st.session_state.current_page == "home":
            home_page()
        elif st.session_state.current_page == "exam_prep":
            exam_prep_page()
        elif st.session_state.current_page == "trivia":
            trivia_page()
        elif st.session_state.current_page == "playing":
            playing_page()
        elif st.session_state.current_page == "leaderboards":
            leaderboards_page()
        elif st.session_state.current_page == "mindfulness":
            mindfulness_page()
        elif st.session_state.current_page == "edit_profile":
            edit_profile_page()
        elif st.session_state.current_page == "lobby_page":
            lobby_page()
    else:
        login_page()

if __name__ == "__main__":
    main()